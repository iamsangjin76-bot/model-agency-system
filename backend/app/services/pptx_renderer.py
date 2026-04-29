# -*- coding: utf-8 -*-
"""
PPTX slide renderer: model data extraction and per-slide layout logic.
"""
from __future__ import annotations

from typing import Any, Optional

from pptx import Presentation
from pptx.util import Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from .pptx_builder import (
    W, H,
    WHITE, DARK, GRAY, LGRAY,
    TMPL, MODEL_TYPE_LABELS, GENDER_LABELS, CAREER_FIELDS,
    _rect, _text, _section_title, _info_row, _fmt_num, _add_photo, _resolve_path,
)


# ---------------------------------------------------------------------------
# Model data extraction
# ---------------------------------------------------------------------------

def _build_data(model: Any) -> dict:
    """Extract a flat dictionary of display-ready strings from a model ORM object."""
    def _str(v: Any) -> str:
        return str(v) if v else ""

    return {
        "name": _str(model.name),
        "name_english": _str(model.name_english),
        "model_type": str(model.model_type.value) if model.model_type else "",
        "gender": str(model.gender.value) if model.gender else "",
        "birth_date": str(model.birth_date) if model.birth_date else "",
        "nationality": _str(model.nationality),
        "school": _str(model.school),
        "debut": _str(model.debut),
        "hobby": _str(model.hobby),
        "languages": _str(model.languages),
        "height": model.height,
        "weight": model.weight,
        "bust": model.bust,
        "waist": model.waist,
        "hip": model.hip,
        "shoe_size": model.shoe_size,
        "agency_name": _str(model.agency_name),
        "agency_phone": _str(model.agency_phone),
        "contact1": _str(model.contact1),
        "instagram_id": _str(model.instagram_id),
        "instagram_followers": model.instagram_followers,
        "youtube_id": _str(model.youtube_id),
        "youtube_subscribers": model.youtube_subscribers,
        "tiktok_id": _str(model.tiktok_id),
        "tiktok_followers": model.tiktok_followers,
        "career_broadcast": _str(model.career_broadcast),
        "career_movie": _str(model.career_movie),
        "career_commercial": _str(model.career_commercial),
        "career_print_ad": _str(model.career_print_ad),
        "career_fashion_show": _str(model.career_fashion_show),
        "career_musical": _str(model.career_musical),
        "career_album": _str(getattr(model, "career_album", "")),
        "career_music_video": _str(getattr(model, "career_music_video", "")),
        "career_other": _str(model.career_other),
        "visa_type": _str(model.visa_type) if hasattr(model, "visa_type") else "",
        "entry_date": str(model.entry_date) if hasattr(model, "entry_date") and model.entry_date else "",
        "memo": _str(model.memo),
        "profile_image": _str(getattr(model, "profile_image", "")),
        # Resolved via files relationship; filled in by build_pptx
        "image_path": None,
    }


# ---------------------------------------------------------------------------
# Slide builder
# ---------------------------------------------------------------------------

def _add_slide(
    prs: Presentation, data: dict, tmpl_key: str,
    upload_dir: str, model_files_dir: str,
) -> None:
    """Add one model slide to the presentation using the given template key."""
    cfg = TMPL.get(tmpl_key, TMPL["new_model_a"])
    color: RGBColor = cfg["color"]
    bg: RGBColor = cfg["bg"]
    layout: str = cfg["layout"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # White background
    _rect(slide, 0, 0, W, H, WHITE)

    # Accent background for sns layout (left info zone)
    if layout == "sns":
        _rect(slide, 0, 0, Cm(18.5), H, bg)

    # Header bar
    _rect(slide, 0, 0, W, Cm(1.5), color)
    header_name = (
        (data.get("name_english") or data.get("name", "")) if layout == "intl"
        else data.get("name", "")
    )
    _text(slide, header_name, Cm(0.5), Cm(0.05), Cm(17), Cm(1.4), 17, True, WHITE)
    type_label = MODEL_TYPE_LABELS.get(data.get("model_type", ""), "")
    _text(slide, type_label, Cm(20.8), Cm(0.05), Cm(8.4), Cm(1.4), 9, False, WHITE, PP_ALIGN.RIGHT)

    # Photo zone — A4 landscape gives more height (21 cm) so the photo can be taller
    ph_w = Cm(8.0)
    ph_h = Cm(14.0)
    ph_t = Cm(2.0)
    ph_l = W - ph_w - Cm(0.5) if layout == "sns" else Cm(0.5)

    img_path = _resolve_path(
        data.get("image_path") or data.get("profile_image"),
        upload_dir, model_files_dir,
    )
    _add_photo(slide, img_path, ph_l, ph_t, ph_w, ph_h, color)

    # Info zone position depends on layout
    if layout == "sns":
        info_l = Cm(0.5)
        info_w = W - ph_w - Cm(1.5)
    else:
        info_l = ph_l + ph_w + Cm(0.7)
        info_w = W - info_l - Cm(0.5)

    y = Cm(2.0)

    # Name block — varies by layout
    if layout == "intl":
        en = data.get("name_english") or ""
        _text(slide, en, info_l, y, info_w, Cm(1.6), 22, True, color)
        y += Cm(1.5)
        _text(slide, data.get("name", ""), info_l, y, info_w, Cm(0.8), 11, False, GRAY)
        y += Cm(0.8)
    elif layout == "bold":
        _text(slide, data.get("name", ""), info_l, y, info_w, Cm(1.8), 24, True, DARK)
        y += Cm(1.6)
        en = data.get("name_english") or ""
        if en:
            _text(slide, en, info_l, y, info_w, Cm(0.75), 10, False, GRAY)
            y += Cm(0.7)
    else:  # classic / sns
        _text(slide, data.get("name", ""), info_l, y, info_w, Cm(1.4), 20, True, DARK)
        y += Cm(1.2)
        en = data.get("name_english") or ""
        if en:
            _text(slide, en, info_l, y, info_w, Cm(0.75), 10, False, GRAY)
            y += Cm(0.7)

    # Measurements pill
    meas = []
    if data.get("height"): meas.append(f"키 {data['height']}cm")
    if data.get("weight"): meas.append(f"몸무게 {data['weight']}kg")
    if data.get("bust"): meas.append(f"가슴 {data['bust']}")
    if data.get("waist"): meas.append(f"허리 {data['waist']}")
    if data.get("hip"): meas.append(f"힙 {data['hip']}")
    if meas:
        _rect(slide, info_l, y, info_w, Cm(0.65), bg)
        _text(slide, "  ·  ".join(meas), info_l + Cm(0.2), y + Cm(0.05),
              info_w - Cm(0.4), Cm(0.55), 9, False, color)
        y += Cm(0.75)
    y += Cm(0.2)

    # SNS stats — prominent in influencer template
    if layout == "sns":
        _section_title(slide, "SNS", info_l, y, color)
        y += Cm(0.65)
        for id_key, cnt_key, prefix in [
            ("instagram_id", "instagram_followers", "@"),
            ("youtube_id", "youtube_subscribers", ""),
            ("tiktok_id", "tiktok_followers", "@"),
        ]:
            val = data.get(id_key) or ""
            if val:
                cnt = _fmt_num(data.get(cnt_key))
                platform = id_key.replace("_id", "").capitalize()
                row = f"{platform}  {prefix}{val}" + (f"  ({cnt})" if cnt else "")
                _text(slide, row, info_l, y, info_w, Cm(0.65), 11, True, color)
                y += Cm(0.65)
        y += Cm(0.2)

    # Basic info section
    _section_title(slide, "기본 정보", info_l, y, color)
    y += Cm(0.65)
    basic: list[tuple[str, str]] = []
    if data.get("gender"): basic.append(("성별", GENDER_LABELS.get(data["gender"], data["gender"])))
    if data.get("birth_date"): basic.append(("생년월일", data["birth_date"]))
    if data.get("nationality"): basic.append(("국적", data["nationality"]))
    if layout == "intl":
        if data.get("visa_type"): basic.append(("비자", data["visa_type"]))
        if data.get("entry_date"): basic.append(("입국일", data["entry_date"]))
        if data.get("languages"): basic.append(("언어", data["languages"]))
    if data.get("school"): basic.append(("출신학교", data["school"]))
    if data.get("debut"): basic.append(("데뷔", data["debut"]))
    for lbl, val in basic:
        _info_row(slide, lbl, val, y, info_l, info_w)
        y += Cm(0.52)
    y += Cm(0.2)

    # Career section
    career = [(lbl, data.get(fld, "")) for lbl, fld in CAREER_FIELDS if data.get(fld)]
    if career and y < H - Cm(4):
        _section_title(slide, "경력", info_l, y, color)
        y += Cm(0.65)
        for lbl, val in career:
            if y > H - Cm(2.5):
                break
            _info_row(slide, lbl, val, y, info_l, info_w)
            y += Cm(0.52)
        y += Cm(0.15)

    # SNS section (non-influencer templates only)
    if layout != "sns" and y < H - Cm(2.5):
        sns_rows: list[tuple[str, str]] = []
        if data.get("instagram_id"):
            cnt = _fmt_num(data.get("instagram_followers"))
            sns_rows.append(("Instagram", f"@{data['instagram_id']}" + (f" ({cnt})" if cnt else "")))
        if data.get("youtube_id"):
            cnt = _fmt_num(data.get("youtube_subscribers"))
            sns_rows.append(("YouTube", data["youtube_id"] + (f" ({cnt})" if cnt else "")))
        if data.get("tiktok_id"):
            cnt = _fmt_num(data.get("tiktok_followers"))
            sns_rows.append(("TikTok", f"@{data['tiktok_id']}" + (f" ({cnt})" if cnt else "")))
        if sns_rows:
            _section_title(slide, "SNS", info_l, y, color)
            y += Cm(0.65)
            for lbl, val in sns_rows:
                if y > H - Cm(2.5):
                    break
                _info_row(slide, lbl, val, y, info_l, info_w)
                y += Cm(0.52)

    # Footer bar
    _rect(slide, 0, H - Cm(0.9), W, Cm(0.9), color)
    _text(slide, "CONFIDENTIAL — 본 자료는 외부 공개를 금합니다",
          Cm(1), H - Cm(0.85), W - Cm(2), Cm(0.8),
          size=7, color=WHITE, align=PP_ALIGN.CENTER)
