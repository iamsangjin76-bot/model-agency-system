# -*- coding: utf-8 -*-
"""
PPTX generation service — public API entry point.
Supports 4 template types: new_model_a, new_model_b, influencer, foreign_model.
"""
from __future__ import annotations

import io
from typing import Any

from pptx import Presentation

from .pptx_builder import W, H, TMPL, _resolve_path
from .pptx_renderer import _build_data, _add_slide


def build_pptx(
    models: list[Any],
    template: str = "new_model_a",
    upload_dir: str = "./uploads",
    model_files_dir: str = "./model_files",
) -> bytes:
    """Build a PPTX bytes blob for the given model ORM objects.

    Args:
        models: List of model ORM instances to render (one slide per model).
        template: Template key — one of new_model_a, new_model_b, influencer, foreign_model.
        upload_dir: Filesystem path for uploaded files (used for image resolution).
        model_files_dir: Filesystem path for model-specific files.

    Returns:
        Raw bytes of the generated .pptx file.
    """
    if template not in TMPL:
        template = "new_model_a"

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    for model in models:
        data = _build_data(model)
        # Resolve profile image from the files relationship
        profile_file = next(
            (f for f in getattr(model, "files", []) if f.is_profile_image), None
        )
        if profile_file:
            data["image_path"] = _resolve_path(
                "/" + profile_file.file_path.lstrip("/"),
                upload_dir, model_files_dir,
            )
        _add_slide(prs, data, template, upload_dir, model_files_dir)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
