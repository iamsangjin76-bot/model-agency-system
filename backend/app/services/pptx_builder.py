# -*- coding: utf-8 -*-
"""
Low-level PPTX primitives: constants, color palette, and helper drawing functions.
All symbols are module-level so they can be imported freely by other modules.
"""
from __future__ import annotations

import logging
import os
from typing import Any, Optional

from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Slide dimensions — A4 landscape
# ---------------------------------------------------------------------------

W = Cm(29.7)
H = Cm(21.0)

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x11, 0x18, 0x27)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LGRAY = RGBColor(0xF3, 0xF4, 0xF6)

# ---------------------------------------------------------------------------
# Template definitions
# ---------------------------------------------------------------------------

TMPL: dict[str, dict] = {
    "new_model_a": dict(
        color=RGBColor(0x7C, 0x3A, 0xED), bg=RGBColor(0xF3, 0xF0, 0xFF), layout="classic",
    ),
    "new_model_b": dict(
        color=RGBColor(0x1F, 0x29, 0x37), bg=RGBColor(0xF9, 0xFA, 0xFB), layout="bold",
    ),
    "influencer": dict(
        color=RGBColor(0xDB, 0x27, 0x77), bg=RGBColor(0xFF, 0xF0, 0xF9), layout="sns",
    ),
    "foreign_model": dict(
        color=RGBColor(0x1D, 0x4E, 0xD8), bg=RGBColor(0xEF, 0xF6, 0xFF), layout="intl",
    ),
}

# ---------------------------------------------------------------------------
# Label lookup tables
# NOTE: Keep these values in sync with frontend/src/types/model.ts
# ---------------------------------------------------------------------------

MODEL_TYPE_LABELS = {
    "new_model": "신인 모델", "influencer": "인플루언서",
    "foreign_model": "외국인 모델", "celebrity": "연예인",
}
GENDER_LABELS = {"male": "남성", "female": "여성", "other": "기타"}

CAREER_FIELDS = [
    ("방송", "career_broadcast"), ("영화", "career_movie"),
    ("광고", "career_commercial"), ("지면광고", "career_print_ad"),
    ("패션쇼", "career_fashion_show"), ("뮤지컬", "career_musical"),
    ("앨범", "career_album"), ("뮤직비디오", "career_music_video"),
    ("기타", "career_other"),
]

# ---------------------------------------------------------------------------
# Primitive drawing helpers
# ---------------------------------------------------------------------------

def _rect(slide: Any, l: Any, t: Any, w: Any, h: Any, color: RGBColor) -> None:
    """Draw a filled rectangle with no border."""
    shape = slide.shapes.add_shape(1, l, t, w, h)  # 1 = MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _text(
    slide: Any, text: str, l: Any, t: Any, w: Any, h: Any,
    size: int = 10, bold: bool = False,
    color: RGBColor = DARK, align: Any = PP_ALIGN.LEFT,
) -> None:
    """Add a textbox with a single run of styled text."""
    if not text:
        return
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _section_title(slide: Any, label: str, l: Any, y: Any, color: RGBColor) -> None:
    """Render an uppercase section heading with a thin rule below it."""
    _text(slide, label.upper(), l, y, Cm(20), Cm(0.5), size=8, bold=True, color=color)
    _rect(slide, l, y + Cm(0.48), Cm(12), Cm(0.04), color)


def _info_row(slide: Any, label: str, value: str, y: Any, l: Any, w: Any) -> None:
    """Render a two-column label/value row."""
    _text(slide, label, l, y, Cm(3.0), Cm(0.5), size=9, color=GRAY)
    _text(slide, value, l + Cm(3.2), y, w - Cm(3.2), Cm(0.5), size=9, color=DARK)


def _fmt_num(n: Optional[int]) -> str:
    """Format a large integer as K/M shorthand; returns empty string for falsy values."""
    if not n:
        return ""
    if n >= 1_000_000:
        return f"{n / 1_000_000:.1f}M"
    if n >= 1_000:
        return f"{n / 1_000:.1f}K"
    return str(n)


def _add_photo(
    slide: Any, image_path: Optional[str],
    l: Any, t: Any, w: Any, h: Any,
    color: RGBColor,
) -> None:
    """Add model photo; draw a placeholder rectangle if the image is missing."""
    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, l, t, w, h)
            return
        except Exception:
            pass
    _rect(slide, l, t, w, h, LGRAY)
    _text(slide, "사진 없음", l, t + h // 2 - Cm(0.3), w, Cm(0.6),
          size=9, color=GRAY, align=PP_ALIGN.CENTER)


def _resolve_path(url: Optional[str], upload_dir: str, model_files_dir: str) -> Optional[str]:
    """Resolve a URL-style path to an absolute filesystem path.

    Security: normalizes the result with os.path.abspath and verifies
    the resolved path is strictly inside the allowed base directory to
    prevent path-traversal attacks (e.g. '../../etc/passwd').
    """
    if not url:
        return None
    for prefix, directory in [("/uploads/", upload_dir), ("/model_files/", model_files_dir)]:
        if url.startswith(prefix):
            base = os.path.abspath(directory)
            candidate = os.path.abspath(os.path.join(directory, url[len(prefix):]))
            # Reject if candidate escapes the base directory
            if not candidate.startswith(base + os.sep) and candidate != base:
                logger.warning("Path traversal attempt blocked: %s", url)
                return None
            if os.path.exists(candidate):
                return candidate
    return None
